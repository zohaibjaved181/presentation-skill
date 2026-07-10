#!/usr/bin/env python3
"""Geometry and visual-coherence lint for PPTX decks."""

from __future__ import annotations

import argparse
import json
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any

import zipfile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from design_tokens import get_style_preset


# Rules that are written against python-pptx's specific geometry conventions
# (margin tokens, accent-rail-under-card pattern, rectangle-card shape type)
# and fire as false positives on pptxgenjs-rendered decks. pptxgenjs places
# content inside the same safe region but uses different shape primitives and
# margin math, so these checks misreport valid layouts as broken.
_PPTXGENJS_SKIPPED_RULES = {
    "margin_left",
    "margin_right",
    "rail_orphan",
    "rail_card_misalignment",
    "rounded_card_with_accent_rail",
}


def _detect_renderer(pptx_path: Path) -> str:
    """Return 'pptxgenjs', 'python-pptx', or 'unknown' based on app.xml."""
    try:
        with zipfile.ZipFile(str(pptx_path), "r") as z:
            with z.open("docProps/app.xml") as f:
                app_xml = f.read().decode("utf-8", errors="replace")
    except (KeyError, zipfile.BadZipFile, OSError):
        return "unknown"
    # pptxgenjs sets Application to "Microsoft Office PowerPoint"; python-pptx
    # leaves the default which LibreOffice rewrites as "Microsoft Macintosh
    # PowerPoint" on render. "PptxGenJS" may also appear in Company.
    if "Microsoft Office PowerPoint" in app_xml or "PptxGenJS" in app_xml:
        return "pptxgenjs"
    if "Microsoft Macintosh PowerPoint" in app_xml or "python-pptx" in app_xml:
        return "python-pptx"
    return "unknown"

# Bug 2: a numeric KPI value may optionally carry a short unit suffix
# (e.g., "14%", "2.1pt", "5x", "98", "$12M", "6 mo"). Reject bare
# adjectives like "Live" / "Higher" / "Clear".
_STATS_NUMERIC_VALUE_RE = re.compile(
    r"^[\$€£¥\-−]?[\d][\d,./\s]*[a-zA-Z%°×x\$€£¥]{0,4}$"
)


def _is_numeric_stats_value(text: str) -> bool:
    stripped = (text or "").strip()
    if not stripped:
        return True  # empty is separately handled upstream
    # Permit common monetary / numeric forms by checking that at least
    # one digit is present before the unit-suffix portion.
    if not any(ch.isdigit() for ch in stripped):
        return False
    return bool(_STATS_NUMERIC_VALUE_RE.match(stripped))


@dataclass
class ShapeInfo:
    shape_id: str
    name: str
    left: float
    top: float
    width: float
    height: float
    has_text: bool
    is_auto_shape: bool
    auto_shape_type_name: str | None
    is_line_like: bool
    area: float

    @property
    def right(self) -> float:
        return self.left + self.width

    @property
    def bottom(self) -> float:
        return self.top + self.height

    @property
    def is_rail_candidate(self) -> bool:
        return self.is_auto_shape and self.width >= 1.2 and self.height <= 0.16

    @property
    def is_card_candidate(self) -> bool:
        return (
            self.is_auto_shape
            and self.area >= 1.0
            and self.width >= 1.4
            and self.height >= 0.7
            and not self.is_line_like
        )

    @property
    def is_rounded_rectangle(self) -> bool:
        return self.auto_shape_type_name == "ROUNDED_RECTANGLE"


def _shape_text(shape: Any) -> str:
    if not hasattr(shape, "text_frame"):
        return ""
    chunks: list[str] = []
    for paragraph in shape.text_frame.paragraphs:
        text = paragraph.text.strip()
        if text:
            chunks.append(text)
    return " ".join(chunks).strip()


def _shape_info(index: int, shape: Any) -> ShapeInfo:
    left = shape.left / 914400.0
    top = shape.top / 914400.0
    width = shape.width / 914400.0
    height = shape.height / 914400.0
    shape_type = shape.shape_type
    is_auto = shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
    auto_shape_type = None
    if is_auto:
        try:
            auto_shape_type = shape.auto_shape_type
        except ValueError:
            auto_shape_type = None
    auto_shape_type_name = (
        getattr(auto_shape_type, "name", str(auto_shape_type)) if auto_shape_type else None
    )
    text = _shape_text(shape)
    return ShapeInfo(
        shape_id=f"shape-{index}",
        name=getattr(shape, "name", f"shape-{index}"),
        left=left,
        top=top,
        width=width,
        height=height,
        has_text=bool(text),
        is_auto_shape=is_auto,
        auto_shape_type_name=auto_shape_type_name,
        is_line_like=(width <= 0.05 or height <= 0.05),
        area=width * height,
    )


def _overlap_width(a: ShapeInfo, b: ShapeInfo) -> float:
    left = max(a.left, b.left)
    right = min(a.right, b.right)
    return max(0.0, right - left)


def _contains_box(outer: ShapeInfo, inner: ShapeInfo, pad: float = 0.02) -> bool:
    return (
        inner.left >= outer.left + pad
        and inner.top >= outer.top + pad
        and inner.right <= outer.right - pad
        and inner.bottom <= outer.bottom - pad
    )


def _cluster(values: list[float], tolerance: float) -> list[list[int]]:
    clusters: list[list[int]] = []
    for idx, value in enumerate(values):
        assigned = False
        for cluster in clusters:
            seed = values[cluster[0]]
            if abs(value - seed) <= tolerance:
                cluster.append(idx)
                assigned = True
                break
        if not assigned:
            clusters.append([idx])
    return clusters


def _group_by_top(cards: list[ShapeInfo], tolerance: float) -> list[list[ShapeInfo]]:
    groups: list[list[ShapeInfo]] = []
    for card in sorted(cards, key=lambda c: c.top):
        placed = False
        for group in groups:
            avg_top = sum(item.top for item in group) / len(group)
            if abs(card.top - avg_top) <= tolerance:
                group.append(card)
                placed = True
                break
        if not placed:
            groups.append([card])
    return groups


def _violation(
    vtype: str,
    shape_ids: list[str],
    severity: str,
    delta: float,
    suggested_fix: str,
    *,
    slide_type: str | None = None,
    extra: dict[str, Any] | None = None,
) -> dict[str, Any]:
    record: dict[str, Any] = {
        "type": vtype,
        "shape_ids": shape_ids,
        "severity": severity,
        "delta_inches": round(delta, 4),
        "suggested_fix": suggested_fix,
    }
    if slide_type:
        record["slide_type"] = slide_type
    if extra:
        record.update(extra)
    return record


def _is_full_bleed_background(shape: ShapeInfo, slide_w: float, slide_h: float) -> bool:
    return (
        shape.left <= 0.03
        and shape.top <= 0.03
        and shape.width >= slide_w - 0.06
        and shape.height >= slide_h - 0.06
    )


# Variant-aware density multipliers. Cards and matrix variants naturally
# cover most of the content zone because they fill it by design; the base
# density cap (tuned for text-heavy layouts) over-flags them. These
# multipliers raise the effective cap for structurally dense variants, while
# keeping a small headroom below 1.0 so true overcrowding can still be flagged.
_VARIANT_DENSITY_MULTIPLIER = {
    "cards-3": 1.35,
    "cards-2": 1.20,
    "matrix": 1.30,
    "timeline": 1.15,
    "stats": 1.22,
    "split": 1.25,
    "generated-image": 1.25,
    # Figure-plus-sidebar evidence slides intentionally fill the main content
    # zone with a large figure panel and interpretation rail.
    "image-sidebar": 1.40,
    "scientific-figure": 1.25,
    "lab-run-results": 1.22,
    "table": 1.24,
    "comparison-2col": 1.30,
}

_WHITESPACE_BALANCE_EXEMPT_VARIANTS = {
    # These variants use deliberate open space as the composition.
    "kpi-hero",
    "pull-quote",
}

_STRUCTURAL_DENSE_VARIANTS = {
    "chart",
    "flow",
    "image-sidebar",
    "lab-run-results",
    "matrix",
    "scientific-figure",
    "stats",
    "table",
}


def _effective_max_density(base: float, outline_slide: dict[str, Any] | None) -> float:
    if not outline_slide:
        return base
    variant = str(outline_slide.get("variant", "")).strip().lower()
    if variant in {"chart", "scientific-figure"}:
        return 1.0
    if variant == "image-sidebar":
        mode = str(outline_slide.get("image_sidebar_mode", "")).strip().lower()
        if mode in {"evidence-mosaic", "editorial-atlas"}:
            return 1.0
    if variant == "comparison-2col":
        return 0.99
    multiplier = _VARIANT_DENSITY_MULTIPLIER.get(variant, 1.0)
    # Structural evidence slides often cover most of the canvas with charts,
    # dashboard panels, native tables, or figure grids. Their readability is
    # governed by text/table/chart checks; the density lint should still catch
    # near-solid overcrowding, but not intentional full-width evidence regions.
    cap = 0.98 if variant in _STRUCTURAL_DENSE_VARIANTS else 0.96
    minimum = 0.97 if variant in _STRUCTURAL_DENSE_VARIANTS else base
    return max(min(cap, base * multiplier), minimum)


def _content_balance_issues(
    shapes: list[ShapeInfo],
    *,
    slide_w: float,
    slide_h: float,
    margin_x: float,
    top_safe: float,
    bottom_safe: float,
    slide_type: str,
    variant: str,
) -> list[dict[str, Any]]:
    """Flag valid-but-awkward slides where content is stranded in a small band.

    Empty-ratio catches broadly sparse slides, but it does not explain whether
    the problem is a dead lower band, a single cramped side, or a genuinely
    intentional open composition. This check is outline-aware and advisory by
    default so authors get a concrete source-level fix without blocking normal
    fast QA unless warning failure is requested.
    """
    if slide_type in {"title", "section"}:
        return []
    if variant in _WHITESPACE_BALANCE_EXEMPT_VARIANTS:
        return []

    region_left = margin_x
    region_top = top_safe
    region_right = slide_w - margin_x
    region_bottom = slide_h - bottom_safe
    region_w = max(0.01, region_right - region_left)
    region_h = max(0.01, region_bottom - region_top)

    boxes: list[tuple[float, float, float, float, str]] = []
    for shape in shapes:
        if shape.area < 0.08 or shape.is_line_like:
            continue
        if _is_full_bleed_background(shape, slide_w, slide_h):
            continue
        if shape.bottom <= region_top + 0.03:
            continue
        if shape.top >= region_bottom - 0.03:
            continue

        left = max(shape.left, region_left)
        top = max(shape.top, region_top)
        right = min(shape.right, region_right)
        bottom = min(shape.bottom, region_bottom)
        if right - left < 0.12 or bottom - top < 0.12:
            continue
        boxes.append((left, top, right, bottom, shape.shape_id))

    if not boxes:
        return []

    union_left = min(item[0] for item in boxes)
    union_top = min(item[1] for item in boxes)
    union_right = max(item[2] for item in boxes)
    union_bottom = max(item[3] for item in boxes)
    span_w_ratio = (union_right - union_left) / region_w
    span_h_ratio = (union_bottom - union_top) / region_h

    left_dead = max(0.0, union_left - region_left)
    right_dead = max(0.0, region_right - union_right)
    top_dead = max(0.0, union_top - region_top)
    bottom_dead = max(0.0, region_bottom - union_bottom)
    max_horizontal_dead_ratio = max(left_dead, right_dead) / region_w
    max_vertical_dead_ratio = max(top_dead, bottom_dead) / region_h
    shape_ids = [item[4] for item in boxes[:4]]

    issues: list[dict[str, Any]] = []
    if span_h_ratio < 0.38 and max_vertical_dead_ratio > 0.42:
        issues.append(
            _violation(
                "content_span_too_short",
                shape_ids,
                "warning",
                0.38 - span_h_ratio,
                (
                    "Content occupies a short vertical band with a large dead zone; "
                    "enlarge the evidence block, add a visual/table/sidebar, or switch "
                    "to a deliberate sparse variant."
                ),
                slide_type=slide_type,
                extra={
                    "variant": variant,
                    "content_span_height_ratio": round(span_h_ratio, 4),
                    "max_vertical_dead_ratio": round(max_vertical_dead_ratio, 4),
                },
            )
        )
    if span_w_ratio < 0.52 and max_horizontal_dead_ratio > 0.42:
        issues.append(
            _violation(
                "content_span_too_narrow",
                shape_ids,
                "warning",
                0.52 - span_w_ratio,
                (
                    "Content is clustered into a narrow side of the safe area; "
                    "use the opposing column for evidence, widen the primary block, "
                    "or choose a one-column sparse variant intentionally."
                ),
                slide_type=slide_type,
                extra={
                    "variant": variant,
                    "content_span_width_ratio": round(span_w_ratio, 4),
                    "max_horizontal_dead_ratio": round(max_horizontal_dead_ratio, 4),
                },
            )
        )
    return issues


def _lint_slide(
    slide: Any,
    slide_index: int,
    slide_w: float,
    slide_h: float,
    *,
    margin_x: float,
    top_safe: float,
    bottom_safe: float,
    gutter: float,
    edge_tol: float,
    gutter_tol: float,
    max_density: float,
    max_empty_ratio: float,
    slide_type: str = "content",
    outline_slide: dict[str, Any] | None = None,
    strict_stats: bool = False,
) -> dict[str, Any]:
    shapes = [_shape_info(i, shape) for i, shape in enumerate(slide.shapes)]
    cards = [
        shape
        for shape in shapes
        if shape.is_card_candidate and not _is_full_bleed_background(shape, slide_w, slide_h)
    ]
    rails = [
        shape
        for shape in shapes
        if shape.is_rail_candidate and not _is_full_bleed_background(shape, slide_w, slide_h)
    ]
    violations: list[dict[str, Any]] = []

    # Margin checks for meaningful blocks.
    for shape in cards:
        if shape.left < margin_x - edge_tol:
            delta = abs(shape.left - margin_x)
            violations.append(
                _violation(
                    "margin_left",
                    [shape.shape_id],
                    "error",
                    delta,
                    "Shift shape right to respect left margin token.",
                )
            )
        if shape.right > (slide_w - margin_x + edge_tol):
            delta = abs(shape.right - (slide_w - margin_x))
            violations.append(
                _violation(
                    "margin_right",
                    [shape.shape_id],
                    "error",
                    delta,
                    "Reduce width or shift left to respect right margin token.",
                )
            )
        if shape.top > top_safe and shape.bottom > (slide_h - bottom_safe + edge_tol):
            delta = abs(shape.bottom - (slide_h - bottom_safe))
            violations.append(
                _violation(
                    "margin_bottom",
                    [shape.shape_id],
                    "warning",
                    delta,
                    "Raise card or reduce height to preserve footer breathing room.",
                )
            )

    variant = str(outline_slide.get("variant", "")).strip().lower() if outline_slide else ""
    evidence_first_variants = {"image-sidebar", "scientific-figure", "lab-run-results", "table", "chart"}
    if variant not in evidence_first_variants:
        text_shapes = [shape for shape in shapes if shape.has_text]
        for card in cards:
            if card.area > slide_w * slide_h * 0.45:
                continue
            inside_text = [
                shape
                for shape in text_shapes
                if shape.shape_id != card.shape_id and _contains_box(card, shape)
            ]
            if not inside_text:
                continue
            fill_ratio = sum(shape.height for shape in inside_text) / max(card.height, 0.01)
            if fill_ratio < 0.18 and card.height >= 1.2:
                violations.append(
                    _violation(
                        "card_text_fill_too_low",
                        [card.shape_id] + [shape.shape_id for shape in inside_text[:3]],
                        "warning",
                        0.18 - fill_ratio,
                        "Use the card area more intentionally or shrink/split the card.",
                        slide_type=slide_type,
                        extra={"text_fill_ratio": round(fill_ratio, 3)},
                    )
                )

    # Row-level consistency checks.
    for group in _group_by_top(cards, tolerance=0.16):
        if len(group) < 2:
            continue
        top_delta = max(item.top for item in group) - min(item.top for item in group)
        if top_delta > edge_tol:
            severity = "error" if top_delta > edge_tol * 2 else "warning"
            violations.append(
                _violation(
                    "top_misalignment",
                    [item.shape_id for item in group],
                    severity,
                    top_delta,
                    "Align top edges for cards in the same visual row.",
                )
            )

        width_delta = max(item.width for item in group) - min(item.width for item in group)
        if width_delta <= 0.5:
            height_delta = max(item.height for item in group) - min(
                item.height for item in group
            )
            if height_delta > edge_tol:
                severity = "error" if height_delta > edge_tol * 2 else "warning"
                violations.append(
                    _violation(
                        "height_inconsistent",
                        [item.shape_id for item in group],
                        severity,
                        height_delta,
                        "Normalize card heights in the same row.",
                    )
                )

        sorted_row = sorted(group, key=lambda item: item.left)
        gaps = []
        for left_shape, right_shape in zip(sorted_row, sorted_row[1:]):
            gap = right_shape.left - left_shape.right
            gaps.append(gap)
            if gap < gutter - gutter_tol:
                violations.append(
                    _violation(
                        "gutter_too_small",
                        [left_shape.shape_id, right_shape.shape_id],
                        "error",
                        gutter - gap,
                        "Increase horizontal spacing to meet gutter token.",
                    )
                )
        if len(gaps) >= 2:
            spread = max(gaps) - min(gaps)
            if spread > 0.25:
                violations.append(
                    _violation(
                        "gutter_inconsistent",
                        [item.shape_id for item in sorted_row],
                        "warning",
                        spread,
                        "Distribute cards evenly for consistent rhythm.",
                    )
                )

    # Column consistency using left-edge clustering.
    if cards:
        left_values = [item.left for item in cards]
        for cluster in _cluster(left_values, tolerance=edge_tol * 2):
            if len(cluster) < 2:
                continue
            members = [cards[idx] for idx in cluster]
            spread = max(member.left for member in members) - min(
                member.left for member in members
            )
            if spread > edge_tol:
                violations.append(
                    _violation(
                        "column_misalignment",
                        [item.shape_id for item in members],
                        "warning",
                        spread,
                        "Align left edges across repeated card columns.",
                    )
                )

    # Rail-to-card alignment checks. Title-slide motifs legitimately use thin
    # decorative stripes without matching card bodies, so skip orphan-rail
    # lint on title slides.
    if slide_type != "title":
        for rail in rails:
            candidates = []
            for card in cards:
                if card.top < rail.top:
                    continue
                if card.top - rail.top > 0.55:
                    continue
                overlap = _overlap_width(rail, card)
                if overlap >= min(rail.width, card.width) * 0.60:
                    candidates.append(card)
            if not candidates:
                violations.append(
                    _violation(
                        "rail_orphan",
                        [rail.shape_id],
                        "warning",
                        rail.height,
                        "Attach accent rail to a card below or remove orphan rail.",
                    )
                )
                continue
            card = min(candidates, key=lambda item: abs(item.top - rail.top))
            left_delta = abs(rail.left - card.left)
            width_delta = abs(rail.width - card.width)
            if left_delta > edge_tol or width_delta > edge_tol:
                severity = "error" if max(left_delta, width_delta) > edge_tol * 2 else "warning"
                violations.append(
                    _violation(
                        "rail_card_misalignment",
                        [rail.shape_id, card.shape_id],
                        severity,
                        max(left_delta, width_delta),
                        "Match rail x/width to corresponding card geometry.",
                    )
                )

            if (
                card.is_rounded_rectangle
                and abs(card.top - rail.top) <= max(0.03, edge_tol)
                and _overlap_width(rail, card) >= min(rail.width, card.width) * 0.9
            ):
                violations.append(
                    _violation(
                        "rounded_card_with_accent_rail",
                        [rail.shape_id, card.shape_id],
                        "error",
                        abs(card.top - rail.top),
                        "Use RECTANGLE for card bodies when applying rectangular accent rails.",
                    )
                )

    # Density model (simple area coverage score).
    content_shapes = [
        shape
        for shape in shapes
        if shape.area >= 0.15 and not shape.is_line_like and shape.top > 0.10
    ]
    covered_area = sum(shape.area for shape in content_shapes)
    overload_shapes = [
        shape
        for shape in content_shapes
        if not _is_full_bleed_background(shape, slide_w, slide_h)
    ]
    overload_area = sum(shape.area for shape in overload_shapes)
    slide_area = max(0.01, slide_w * slide_h)
    density = min(1.0, covered_area / slide_area)
    overload_density = min(1.0, overload_area / slide_area)
    empty_ratio = max(0.0, 1.0 - density)

    effective_max_density = _effective_max_density(max_density, outline_slide)
    if slide_type != "title" and overload_density > effective_max_density:
        violations.append(
            _violation(
                "density_too_high",
                [],
                "warning",
                overload_density - effective_max_density,
                "Reduce number/size of blocks or split content across slides.",
                slide_type=slide_type,
                extra={"visual_density_score": round(density, 4)},
            )
        )

    # Bug 5 follow-up: section dividers should be visually fuller than title
    # openers, but title slides are intentionally sparse when they rely on the
    # default motif. Keep section slides strict; make title-slide empty ratio
    # advisory unless it becomes truly extreme.
    effective_max_empty_ratio = max_empty_ratio
    if slide_type == "title":
        effective_max_empty_ratio = max(max_empty_ratio, 0.76)
    elif slide_type == "section":
        effective_max_empty_ratio = max(max_empty_ratio * 0.80, 0.50)

    # Sparse-by-design variants (kpi-hero, comparison-2col, pull-quote) use
    # deliberate whitespace as visual emphasis; raise the empty-ratio cap
    # so the lint doesn't fight the composition.
    if outline_slide:
        _sparse_variant_caps = {
            "kpi-hero": 0.85,
            "pull-quote": 0.85,
            "comparison-2col": 0.85,
            "split": 0.66,
            "stats": 0.66,
            "timeline": 0.78,
            "matrix": 0.72,
        }
        _variant = str(outline_slide.get("variant", "")).strip().lower()
        if _variant in _sparse_variant_caps:
            effective_max_empty_ratio = max(
                effective_max_empty_ratio,
                _sparse_variant_caps[_variant],
            )

    if empty_ratio > effective_max_empty_ratio:
        delta = empty_ratio - effective_max_empty_ratio
        severity = "warning"
        if slide_type == "section" and delta > 0.10:
            severity = "error"
        violations.append(
            _violation(
                "empty_ratio_too_high",
                [],
                severity,
                delta,
                "Increase purposeful visual structure to reduce dead zones.",
                slide_type=slide_type,
                extra={"effective_max_empty_ratio": round(effective_max_empty_ratio, 4)},
            )
        )

    violations.extend(
        _content_balance_issues(
            shapes,
            slide_w=slide_w,
            slide_h=slide_h,
            margin_x=margin_x,
            top_safe=top_safe,
            bottom_safe=bottom_safe,
            slide_type=slide_type,
            variant=variant,
        )
    )

    # Bug 2: flag non-numeric stats KPI values. The outline tells us
    # which slides are stats variants and what the declared `value`
    # strings are. We warn (default) or error (with --strict-stats) so
    # authors can swap adjectives like "Higher" for a numeric or a
    # clearly qualitative variant.
    if outline_slide and _is_stats_slide(outline_slide):
        stats_issues = _stats_value_issues(
            outline_slide, slide_type=slide_type, strict_stats=strict_stats
        )
        violations.extend(stats_issues)

    # Fix 1: validate inline chart schemas so malformed chart payloads are
    # caught even when the renderer substituted a red error banner. This is
    # a blocking error because it signals the slide's primary visual is
    # missing.
    if outline_slide and _is_chart_slide(outline_slide):
        chart_issues = _chart_schema_issues(outline_slide, slide_type=slide_type)
        violations.extend(chart_issues)

    return {
        "slide_index": slide_index,
        "slide_type": slide_type,
        "violations": violations,
        "density_score": round(density, 4),
        "empty_ratio": round(empty_ratio, 4),
    }


def _is_chart_slide(outline_slide: dict[str, Any]) -> bool:
    variant = str(outline_slide.get("variant", "")).strip().lower()
    if variant == "chart":
        return True
    # Fallback: slides without an explicit variant but with a `chart` object
    # also render via the chart variant in build_deck.py.
    chart = outline_slide.get("chart")
    assets = outline_slide.get("assets") if isinstance(outline_slide.get("assets"), dict) else {}
    return bool(chart) or bool(assets.get("chart_data") or assets.get("chart"))


def _chart_schema_issues(
    outline_slide: dict[str, Any],
    *,
    slide_type: str,
) -> list[dict[str, Any]]:
    """Validate inline chart payloads in the outline.

    Accepts both accepted schema forms:
      - legacy: series[i] carries its own `labels` list
      - shorthand: chart-level `categories` with series carrying only `values`

    Emits `chart_schema_invalid` (severity error) with a concrete reason so
    QA flags the malformed chart even if the renderer drew a red banner.
    """
    issues: list[dict[str, Any]] = []
    chart = outline_slide.get("chart")
    # Only inline dicts are validated here; staged chart:<alias> JSON files
    # are resolved at build time and fall outside the outline schema.
    if not isinstance(chart, dict):
        return issues

    reasons: list[str] = []
    chart_categories = chart.get("categories")
    if not (isinstance(chart_categories, list) and chart_categories):
        chart_categories = chart.get("labels") if isinstance(chart.get("labels"), list) else None

    series_items = chart.get("series")
    if isinstance(series_items, list) and series_items:
        for index, item in enumerate(series_items):
            if not isinstance(item, dict):
                reasons.append(f"series[{index}] is not an object")
                continue
            series_labels = item.get("labels") if isinstance(item.get("labels"), list) else None
            labels = series_labels if series_labels else chart_categories
            values = item.get("values")
            if not (isinstance(labels, list) and labels):
                reasons.append(
                    f"series[{index}] missing labels (no series.labels and no chart.categories)"
                )
                continue
            if not (isinstance(values, list) and values):
                reasons.append(f"series[{index}] missing values")
                continue
            if len(labels) != len(values):
                reasons.append(
                    f"series[{index}] length mismatch: {len(labels)} labels vs {len(values)} values"
                )
                continue
            bad_vals = [
                v for v in values if not isinstance(v, (int, float)) or isinstance(v, bool)
            ]
            if bad_vals:
                reasons.append(f"series[{index}] contains non-numeric values")
    else:
        # Flat shorthand: {labels|categories: [...], values: [...]}
        values = chart.get("values")
        if not (isinstance(chart_categories, list) and chart_categories):
            reasons.append("chart has no series and no categories/labels")
        elif not (isinstance(values, list) and values):
            reasons.append("chart has no series and no flat values array")
        elif len(chart_categories) != len(values):
            reasons.append(
                f"flat chart length mismatch: {len(chart_categories)} labels vs {len(values)} values"
            )

    if reasons:
        issues.append(
            _violation(
                "chart_schema_invalid",
                [],
                "error",
                0.0,
                (
                    "Chart data must be either {categories:[...], series:[{name, values}]} "
                    "or {series:[{name, labels, values}]}; labels and values must be the "
                    "same length and values must be numeric."
                ),
                slide_type=slide_type,
                extra={"reasons": reasons},
            )
        )
    return issues


def _is_stats_slide(outline_slide: dict[str, Any]) -> bool:
    variant = str(outline_slide.get("variant", "")).strip().lower()
    if variant == "stats":
        return True
    # Fallback: a slide that has `facts` but no explicit variant still
    # renders as stats in build_deck.py.
    facts = outline_slide.get("facts")
    return isinstance(facts, list) and bool(facts)


def _stats_value_issues(
    outline_slide: dict[str, Any],
    *,
    slide_type: str,
    strict_stats: bool,
) -> list[dict[str, Any]]:
    issues: list[dict[str, Any]] = []
    facts = outline_slide.get("facts")
    if not isinstance(facts, list):
        return issues
    severity = "error" if strict_stats else "warning"
    for index, fact in enumerate(facts):
        if not isinstance(fact, dict):
            continue
        value = str(fact.get("value", "")).strip()
        if not value:
            continue
        if _is_numeric_stats_value(value):
            continue
        issues.append(
            _violation(
                "stats_value_non_numeric",
                [],
                severity,
                0.0,
                (
                    "Stats slides render `value` at KPI size; replace the "
                    "adjective with a number, a number+unit, or switch to "
                    "a different variant."
                ),
                slide_type=slide_type,
                extra={
                    "fact_index": index,
                    "value": value,
                    "label": str(fact.get("label", "")).strip(),
                },
            )
        )
    return issues


def _args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Run geometry and density lint on a PPTX deck.")
    parser.add_argument("--input", required=True, help="Input PPTX path")
    parser.add_argument("--output", help="Optional output JSON path")
    parser.add_argument(
        "--style-preset",
        default="executive-clinical",
        help="Style preset token source",
    )
    parser.add_argument("--edge-tolerance", type=float, help="Override edge tolerance in inches")
    parser.add_argument("--gutter-tolerance", type=float, help="Override gutter tolerance")
    parser.add_argument("--max-density", type=float, help="Override max density threshold")
    parser.add_argument("--max-empty-ratio", type=float, help="Override max empty-ratio threshold")
    parser.add_argument(
        "--fail-on-error",
        action="store_true",
        help="Return non-zero if any error-severity violations are found",
    )
    parser.add_argument(
        "--outline",
        help=(
            "Optional path to the outline JSON used to build the deck. "
            "When supplied, slide-type-aware checks (empty_ratio on title/"
            "section, stats_value_non_numeric) become available."
        ),
    )
    parser.add_argument(
        "--strict-stats",
        action="store_true",
        help=(
            "Escalate stats_value_non_numeric findings from warning to error "
            "(Bug 2 in presentation-skill renderer fixes)."
        ),
    )
    return parser.parse_args()


def _load_outline_types(outline_path: Path) -> list[dict[str, Any]]:
    try:
        raw = json.loads(outline_path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return []
    slides = raw.get("slides")
    if not isinstance(slides, list):
        return []
    return [item if isinstance(item, dict) else {} for item in slides]


def main() -> int:
    args = _args()
    input_path = Path(args.input).expanduser().resolve()
    if not input_path.exists():
        raise FileNotFoundError(f"Input not found: {input_path}")

    preset = get_style_preset(args.style_preset)
    layout = preset.layout
    edge_tol = args.edge_tolerance if args.edge_tolerance is not None else layout.edge_tolerance
    gutter_tol = (
        args.gutter_tolerance
        if args.gutter_tolerance is not None
        else layout.gutter_tolerance
    )
    max_density = args.max_density if args.max_density is not None else layout.max_density
    max_empty_ratio = (
        args.max_empty_ratio if args.max_empty_ratio is not None else layout.max_empty_ratio
    )

    prs = Presentation(str(input_path))
    slide_w = prs.slide_width / 914400.0
    slide_h = prs.slide_height / 914400.0

    outline_slides: list[dict[str, Any]] = []
    if args.outline:
        outline_slides = _load_outline_types(Path(args.outline).expanduser().resolve())

    slides_report = []
    for slide_index, slide in enumerate(prs.slides):
        outline_slide = (
            outline_slides[slide_index]
            if slide_index < len(outline_slides)
            else None
        )
        slide_type = str(
            (outline_slide or {}).get("type", "content")
        ).strip().lower() or "content"
        slides_report.append(
            _lint_slide(
                slide,
                slide_index=slide_index,
                slide_w=slide_w,
                slide_h=slide_h,
                margin_x=layout.margin_x,
                top_safe=layout.top_safe,
                bottom_safe=layout.bottom_safe,
                gutter=layout.gutter,
                edge_tol=edge_tol,
                gutter_tol=gutter_tol,
                max_density=max_density,
                max_empty_ratio=max_empty_ratio,
                slide_type=slide_type,
                outline_slide=outline_slide,
                strict_stats=args.strict_stats,
            )
        )

    renderer = _detect_renderer(input_path)
    if renderer == "pptxgenjs":
        # Drop rules that misreport on pptxgenjs output; keep overflow,
        # overlap, empty-ratio, density, and schema checks which are
        # renderer-agnostic.
        for slide in slides_report:
            slide["violations"] = [
                v for v in slide.get("violations", [])
                if v.get("type") not in _PPTXGENJS_SKIPPED_RULES
            ]

    violations = [
        violation for slide in slides_report for violation in slide.get("violations", [])
    ]
    error_count = sum(1 for item in violations if item.get("severity") == "error")
    warning_count = sum(1 for item in violations if item.get("severity") == "warning")
    payload = {
        "input": str(input_path),
        "style_preset": preset.name,
        "renderer_detected": renderer,
        "summary": {
            "slide_count": len(slides_report),
            "violation_count": len(violations),
            "error_count": error_count,
            "warning_count": warning_count,
            "density_score_by_slide": [
                {"slide_index": slide["slide_index"], "density_score": slide["density_score"]}
                for slide in slides_report
            ],
        },
        "slides": slides_report,
    }

    if args.output:
        output_path = Path(args.output).expanduser().resolve()
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_text(json.dumps(payload, indent=2), encoding="utf-8")
        print(f"Layout lint report: {output_path}")
    else:
        print(json.dumps(payload, indent=2))

    if args.fail_on_error and error_count > 0:
        return 1
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
