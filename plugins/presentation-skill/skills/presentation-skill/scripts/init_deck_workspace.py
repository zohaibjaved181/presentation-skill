#!/usr/bin/env python3
"""Create a persistent workspace for iterative deck authoring."""

from __future__ import annotations

import argparse
import json
import re
import struct
import sys
import zlib
from pathlib import Path
from typing import Any

from pptx import Presentation

from design_tokens import PRESETS
from emit_deck_start_packet import build_packet
from model_adaptive_workflow import PROFILE_ALIASES, write_agent_brief
from style_reference_catalog import LAYOUT_PLAYBOOK_VERSION, preset_style_reference
from style_treatment_profiles import preset_treatment_profile
from workflow_atom_context import build_workflow_atom_context, compact_workflow_atom_context


STYLE_REFERENCE_ASSETS = {
    "figure_a": "assets/style_reference/starter_figure_a.png",
    "figure_b": "assets/style_reference/starter_figure_b.png",
    "figure_c": "assets/style_reference/starter_figure_c.png",
    "flow": "assets/style_reference/starter_flow.mmd",
    "icon": "assets/icons/starter_icon.png",
}


def _slugify(value: str) -> str:
    text = re.sub(r"[^A-Za-z0-9]+", "-", value.strip().lower()).strip("-")
    return text or "deck-workspace"


def _copy_json(src: Path) -> dict[str, Any]:
    return json.loads(src.read_text(encoding="utf-8"))


def _shape_text(shape: Any) -> str:
    if not hasattr(shape, "text_frame"):
        return ""
    chunks: list[str] = []
    for paragraph in shape.text_frame.paragraphs:
        text = paragraph.text.strip()
        if text:
            chunks.append(text)
    return "\n".join(chunks).strip()


def _slide_title(slide: Any) -> str:
    try:
        title_shape = slide.shapes.title
        if title_shape is not None:
            return (title_shape.text or "").strip()
    except Exception:
        pass
    return ""


def _extract_outline(reference_pptx: Path) -> dict[str, Any]:
    prs = Presentation(str(reference_pptx))
    slides: list[dict[str, Any]] = []
    for index, slide in enumerate(prs.slides, start=1):
        title = _slide_title(slide)
        body_lines: list[str] = []
        for shape in slide.shapes:
            text = _shape_text(shape)
            if not text:
                continue
            if title and text.strip() == title:
                continue
            for line in text.splitlines():
                cleaned = line.strip()
                if cleaned:
                    body_lines.append(cleaned)
        if index == 1:
            subtitle = body_lines[0] if body_lines else ""
            slides.append({"type": "title", "title": title or "Presentation", "subtitle": subtitle})
            continue
        slide_spec: dict[str, Any] = {
            "type": "content",
            "title": title or f"Slide {index}",
        }
        if body_lines:
            slide_spec["bullets"] = body_lines[:6]
        slides.append(slide_spec)
    return {"slides": slides}


def _ensure_outline_slide_ids(outline: dict[str, Any]) -> list[dict[str, str]]:
    """Add stable slide IDs where missing and return IDs/variants for planning stubs."""
    refs: list[dict[str, str]] = []
    slides = outline.get("slides")
    if not isinstance(slides, list):
        return refs
    for index, slide in enumerate(slides, start=1):
        if not isinstance(slide, dict):
            continue
        slide_id = ""
        for key in ("slide_id", "id", "slug"):
            text = str(slide.get(key) or "").strip()
            if text:
                slide_id = text
                break
        if not slide_id:
            slide_id = f"s{index}"
            slide["slide_id"] = slide_id
        variant = str(slide.get("variant") or "").strip()
        slide_type = str(slide.get("type") or "").strip()
        if not variant and slide_type and slide_type != "content":
            variant = slide_type
        refs.append(
            {
                "slide_id": slide_id,
                "variant": variant,
                "title": str(slide.get("title") or "").strip(),
                "starter_kind": str(slide.get("starter_kind") or "").strip(),
                "treatment_key": str(slide.get("treatment_key") or "").strip(),
            }
        )
    return refs


def _font_families(reference_pptx: Path) -> list[str]:
    prs = Presentation(str(reference_pptx))
    families: set[str] = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            for paragraph in shape.text_frame.paragraphs:
                paragraph_font = getattr(paragraph, "font", None)
                if paragraph_font is not None:
                    name = getattr(paragraph_font, "name", None)
                    if name:
                        families.add(name)
                for run in paragraph.runs:
                    name = getattr(run.font, "name", None)
                    if name:
                        families.add(name)
    return sorted(families)


def _reference_summary(reference_pptx: Path) -> dict[str, Any]:
    prs = Presentation(str(reference_pptx))
    return {
        "reference_pptx": str(reference_pptx),
        "slide_count": len(prs.slides),
        "slide_size_inches": {
            "width": round(prs.slide_width / 914400.0, 3),
            "height": round(prs.slide_height / 914400.0, 3),
        },
        "font_families": _font_families(reference_pptx),
    }


def _story_list(value: Any) -> list[Any]:
    return value if isinstance(value, list) else []


def _story_dict(value: Any) -> dict[str, Any]:
    return value if isinstance(value, dict) else {}


def _compact_text(value: Any, fallback: str, *, limit: int = 64) -> str:
    text = str(value or "").strip() or fallback
    if len(text) <= limit:
        return text
    return text[: max(0, limit - 3)].rstrip(" -;:,") + "..."


def _starter_base_slide(
    *,
    slide_id: str,
    variant: str,
    title: str,
    subtitle: str,
    reference: dict[str, Any],
    treatment_key: str,
) -> dict[str, Any]:
    ref_id = str(reference.get("reference_id") or "style-reference").strip()
    short_ref = ref_id[4:] if ref_id.startswith("ref-") else ref_id
    return {
        "slide_id": slide_id,
        "type": "content",
        "variant": variant,
        "starter_kind": "style_reference",
        "treatment_key": treatment_key,
        "title": title,
        "subtitle": subtitle,
        "footer": "style-reference starter",
        "sources": ["Synthetic style scaffold"],
        "refs": [_compact_text(short_ref, "style-reference", limit=18)],
    }


def _chart_starter_slide(slide_id: str, reference: dict[str, Any]) -> dict[str, Any] | None:
    story = _story_dict(reference.get("example_storyboard"))
    chart = _story_dict(story.get("chart"))
    labels = [str(item) for item in _story_list(chart.get("labels")) if str(item).strip()]
    values = [
        item
        for item in _story_list(chart.get("values"))
        if isinstance(item, (int, float)) and not isinstance(item, bool)
    ]
    if not labels or len(labels) != len(values):
        return None
    treatments = _story_dict(reference.get("content_treatments"))
    title = _compact_text(chart.get("title"), f"{reference.get('reference_name', 'Reference')} chart", limit=58)
    slide = _starter_base_slide(
        slide_id=slide_id,
        variant="chart",
        title="Chart Scaffold",
        subtitle=str(treatments.get("chart") or "Replace with the deck's real chart evidence."),
        reference=reference,
        treatment_key="chart",
    )
    slide["chart_treatment"] = "standard"
    slide["chart"] = {
        "type": "bar",
        "title": title,
        "labels": labels,
        "values": values,
        "notes": str(chart.get("note") or treatments.get("chart") or ""),
        "options": {"showLegend": False, "catAxisLabelFontSize": 8, "valAxisLabelFontSize": 8},
    }
    return slide


def _table_starter_slide(slide_id: str, reference: dict[str, Any], *, variant: str = "table") -> dict[str, Any] | None:
    story = _story_dict(reference.get("example_storyboard"))
    table = _story_dict(story.get("table"))
    headers = [str(item) for item in _story_list(table.get("headers")) if str(item).strip()]
    rows = [
        [str(cell) for cell in row]
        for row in _story_list(table.get("rows"))
        if isinstance(row, list)
    ]
    if not headers or not rows:
        return None
    treatments = _story_dict(reference.get("content_treatments"))
    slide = _starter_base_slide(
        slide_id=slide_id,
        variant="lab-run-results" if variant == "lab-run-results" else "table",
        title="Table Scaffold",
        subtitle=str(treatments.get("table") or "Replace rows with real evidence before delivery."),
        reference=reference,
        treatment_key="table",
    )
    slide["headers"] = headers
    slide["rows"] = rows[:5]
    slide["table_treatment"] = "compact-ledger"
    slide["caption"] = "Synthetic style-reference rows; replace with source-backed data."
    return slide


def _comparison_starter_slide(slide_id: str, reference: dict[str, Any], *, variant: str = "comparison-2col") -> dict[str, Any] | None:
    story = _story_dict(reference.get("example_storyboard"))
    comparison = _story_dict(story.get("comparison"))
    left_body = [str(item) for item in _story_list(comparison.get("left_body")) if str(item).strip()]
    right_body = [str(item) for item in _story_list(comparison.get("right_body")) if str(item).strip()]
    if not left_body and not right_body:
        return None
    treatments = _story_dict(reference.get("content_treatments"))
    if variant == "split":
        slide = _starter_base_slide(
            slide_id=slide_id,
            variant="split",
            title="Contrast Scaffold",
            subtitle=str(treatments.get("comparison") or "Replace with the deck's real contrast."),
            reference=reference,
            treatment_key="comparison",
        )
        slide["bullets"] = left_body[:2] + right_body[:2]
        slide["highlights_label"] = str(comparison.get("right_title") or "Reference move")
        slide["highlights"] = [str(comparison.get("verdict") or "Use the contrast to make the decision visible.")]
        return slide
    slide = _starter_base_slide(
        slide_id=slide_id,
        variant="comparison-2col",
        title="Comparison Scaffold",
        subtitle=str(treatments.get("comparison") or "Replace with the deck's real comparison."),
        reference=reference,
        treatment_key="comparison",
    )
    slide["left"] = {"title": str(comparison.get("left_title") or "Before"), "body": left_body or ["Current state"]}
    slide["right"] = {"title": str(comparison.get("right_title") or "After"), "body": right_body or ["Target state"]}
    slide["verdict"] = str(comparison.get("verdict") or "Make the chosen state explicit.")
    return slide


def _dashboard_starter_slide(slide_id: str, reference: dict[str, Any], *, variant: str = "stats") -> dict[str, Any] | None:
    story = _story_dict(reference.get("example_storyboard"))
    treatments = _story_dict(reference.get("content_treatments"))
    if variant == "kpi-hero":
        kpi = _story_dict(story.get("kpi"))
        if not kpi.get("value") or not kpi.get("label"):
            return None
        slide = _starter_base_slide(
            slide_id=slide_id,
            variant="kpi-hero",
            title="Hero Metric Scaffold",
            subtitle=str(treatments.get("dashboard") or "Use only when one metric carries the point."),
            reference=reference,
            treatment_key="dashboard",
        )
        slide["value"] = str(kpi.get("value"))
        slide["label"] = str(kpi.get("label"))
        slide["context"] = str(kpi.get("context") or "Replace with a source-backed metric.")
        return slide
    facts = [
        {
            "value": str(item.get("value") or ""),
            "label": str(item.get("label") or ""),
            "detail": str(item.get("detail") or ""),
        }
        for item in _story_list(story.get("dashboard_facts"))
        if isinstance(item, dict)
    ]
    facts = [item for item in facts if item["value"] or item["label"]]
    if not facts:
        return None
    slide = _starter_base_slide(
        slide_id=slide_id,
        variant="stats",
        title="Dashboard Scaffold",
        subtitle=str(treatments.get("dashboard") or "Replace with real operating facts."),
        reference=reference,
        treatment_key="dashboard",
    )
    slide["facts"] = facts[:3]
    slide["assets"] = {"icons": [STYLE_REFERENCE_ASSETS["icon"]] * min(3, len(facts))}
    return slide


def _timeline_starter_slide(slide_id: str, reference: dict[str, Any]) -> dict[str, Any] | None:
    story = _story_dict(reference.get("example_storyboard"))
    milestones = [
        {
            "label": str(item.get("label") or ""),
            "title": str(item.get("title") or ""),
            "body": str(item.get("body") or ""),
        }
        for item in _story_list(story.get("timeline"))
        if isinstance(item, dict)
    ]
    milestones = [item for item in milestones if item["title"] or item["body"]]
    if not milestones:
        return None
    slide = _starter_base_slide(
        slide_id=slide_id,
        variant="timeline",
        title="Sequence Scaffold",
        subtitle="Use this only when sequence changes the decision.",
        reference=reference,
        treatment_key="decision",
    )
    slide["milestones"] = milestones[:4]
    return slide


def _matrix_starter_slide(slide_id: str, reference: dict[str, Any]) -> dict[str, Any] | None:
    story = _story_dict(reference.get("example_storyboard"))
    quadrants = [
        {"title": str(item.get("title") or ""), "body": str(item.get("body") or "")}
        for item in _story_list(story.get("quadrants"))
        if isinstance(item, dict)
    ]
    quadrants = [item for item in quadrants if item["title"] or item["body"]]
    if len(quadrants) < 2:
        return None
    slide = _starter_base_slide(
        slide_id=slide_id,
        variant="matrix",
        title="Tradeoff Scaffold",
        subtitle=str(_story_dict(reference.get("content_treatments")).get("comparison") or "Replace with a real decision matrix."),
        reference=reference,
        treatment_key="comparison",
    )
    slide["quadrants"] = quadrants[:4]
    slide["assets"] = {"icons": [STYLE_REFERENCE_ASSETS["icon"]] * min(4, len(slide["quadrants"]))}
    return slide


def _figure_sidebar_sections(reference: dict[str, Any]) -> list[dict[str, Any]]:
    story = _story_dict(reference.get("example_storyboard"))
    figure = _story_dict(story.get("figure"))
    sections: list[dict[str, Any]] = []
    for item in _story_list(figure.get("sections")):
        if not isinstance(item, dict):
            continue
        title = _compact_text(item.get("title"), "Readout", limit=20)
        lines = [
            _compact_text(line, "Replace with figure evidence.", limit=52)
            for line in _story_list(item.get("body"))
            if str(line).strip()
        ]
        if title and lines:
            sections.append({"title": title, "body": lines[:2]})
    if sections:
        return sections[:3]
    return [
        {"title": "Readout", "body": ["One local figure owns the proof."]},
        {"title": "Interpretation", "body": ["Sidebar explains why it matters."]},
    ]


def _figure_starter_slide(slide_id: str, reference: dict[str, Any], *, variant: str) -> dict[str, Any] | None:
    story = _story_dict(reference.get("example_storyboard"))
    figure = _story_dict(story.get("figure"))
    treatments = _story_dict(reference.get("content_treatments"))
    caption = _compact_text(
        figure.get("caption"),
        "Synthetic starter figure; replace with source-backed evidence.",
        limit=86,
    )
    interpretation = _compact_text(
        figure.get("interpretation"),
        "Use this figure grammar as style memory, then replace the evidence.",
        limit=92,
    )
    if variant == "flow":
        slide = _starter_base_slide(
            slide_id=slide_id,
            variant="flow",
            title="Flow Scaffold",
            subtitle=str(treatments.get("figure") or "Replace with a real process diagram."),
            reference=reference,
            treatment_key="figure",
        )
        slide["assets"] = {"mermaid_source": STYLE_REFERENCE_ASSETS["flow"]}
        slide["sidebar_sections"] = _figure_sidebar_sections(reference)[:2]
        slide["summary_callout"] = interpretation
        slide["caption"] = caption
        return slide
    if variant == "scientific-figure":
        slide = _starter_base_slide(
            slide_id=slide_id,
            variant="scientific-figure",
            title="Figure Panels Scaffold",
            subtitle=str(treatments.get("figure") or "Replace with real figure panels."),
            reference=reference,
            treatment_key="figure",
        )
        slide["figures"] = [
            {"path": STYLE_REFERENCE_ASSETS["figure_a"], "label": "A", "title": "Signal", "caption": "Starter panel"},
            {"path": STYLE_REFERENCE_ASSETS["figure_b"], "label": "B", "title": "Contrast", "caption": "Starter panel"},
        ]
        slide["caption"] = caption
        slide["interpretation"] = interpretation
        return slide
    slide = _starter_base_slide(
        slide_id=slide_id,
        variant="image-sidebar",
        title="Figure Sidebar Scaffold",
        subtitle=str(treatments.get("figure") or "Replace with a real figure and interpretation."),
        reference=reference,
        treatment_key="figure",
    )
    slide["assets"] = {"hero_image": STYLE_REFERENCE_ASSETS["figure_a"]}
    slide["image_side"] = "left"
    slide["sidebar_sections"] = _figure_sidebar_sections(reference)
    slide["caption"] = caption
    return slide


def _cards_starter_slide(slide_id: str, reference: dict[str, Any], *, variant: str = "cards-3") -> dict[str, Any] | None:
    treatments = _story_dict(reference.get("content_treatments"))
    cards = [
        {
            "title": "Claim",
            "body": _compact_text(treatments.get("title"), "One clear claim opens the section.", limit=76),
            "accent": "accent_primary",
        },
        {
            "title": "Proof",
            "body": _compact_text(treatments.get("chart"), "Attach one evidence object.", limit=76),
            "accent": "accent_secondary",
        },
        {
            "title": "Decision",
            "body": _compact_text(treatments.get("decision"), "State the implication.", limit=76),
            "accent": "accent_primary",
        },
    ]
    slide = _starter_base_slide(
        slide_id=slide_id,
        variant="cards-2" if variant == "cards-2" else "cards-3",
        title="Proof-Block Scaffold",
        subtitle="Use cards only when the content is genuinely parallel.",
        reference=reference,
        treatment_key="dashboard",
    )
    slide["cards"] = cards[:2] if variant == "cards-2" else cards
    slide["assets"] = {"icons": [STYLE_REFERENCE_ASSETS["icon"]] * len(slide["cards"])}
    if variant != "cards-2":
        slide["promote_card"] = 0
    return slide


def _starter_slide_for_variant(slide_id: str, variant: str, reference: dict[str, Any]) -> dict[str, Any] | None:
    normalized = str(variant or "").strip().lower()
    if normalized == "chart":
        return _chart_starter_slide(slide_id, reference)
    if normalized in {"table", "lab-run-results"}:
        return _table_starter_slide(slide_id, reference, variant=normalized)
    if normalized in {"comparison-2col", "split"}:
        return _comparison_starter_slide(slide_id, reference, variant=normalized)
    if normalized in {"stats", "kpi-hero"}:
        return _dashboard_starter_slide(slide_id, reference, variant=normalized)
    if normalized == "timeline":
        return _timeline_starter_slide(slide_id, reference)
    if normalized == "matrix":
        return _matrix_starter_slide(slide_id, reference)
    if normalized in {"cards-2", "cards-3"}:
        return _cards_starter_slide(slide_id, reference, variant=normalized)
    if normalized in {"image-sidebar", "scientific-figure", "flow"}:
        return _figure_starter_slide(slide_id, reference, variant=normalized)
    return None


def _style_reference_starter_slides(style_preset: str, *, start_index: int = 3, max_slides: int = 3) -> list[dict[str, Any]]:
    reference = preset_style_reference(style_preset)
    playbook = _story_dict(reference.get("layout_playbook"))
    showcase = [
        str(item).strip().lower()
        for item in _story_list(playbook.get("gallery_showcase_variants"))
        if str(item).strip()
    ]
    preferred = [
        str(item).strip().lower()
        for item in _story_list(playbook.get("preferred_variants"))
        if str(item).strip()
    ]
    fallbacks = ["chart", "table", "comparison-2col", "stats", "timeline", "matrix", "kpi-hero", "cards-3"]
    candidates: list[str] = []
    for variant in [*showcase, *preferred, *fallbacks]:
        if variant not in candidates:
            candidates.append(variant)
    slides: list[dict[str, Any]] = []
    treatment_keys: set[str] = set()
    skip_default_starter_variants = {
        "generated-image",
        "timeline",
    }
    skip_by_preset = {
        "sunset-investor": {"kpi-hero"},
        "warm-terracotta": {"cards-2", "split"},
    }
    for variant in candidates:
        if variant in skip_default_starter_variants or variant in skip_by_preset.get(style_preset, set()):
            continue
        slide_id = f"s{start_index + len(slides)}"
        slide = _starter_slide_for_variant(slide_id, variant, reference)
        if not slide:
            continue
        key = str(slide.get("treatment_key") or variant)
        if key in treatment_keys and len(slides) >= 2:
            continue
        treatment_keys.add(key)
        slides.append(slide)
        if len(slides) >= max_slides:
            break
    return slides


def _starter_outline(
    title: str,
    style_preset: str,
    font_pair: str | None,
    palette_key: str | None,
    *,
    user_prompt: str = "",
) -> dict[str, Any]:
    reference = preset_style_reference(style_preset)
    playbook = _story_dict(reference.get("layout_playbook"))
    atom_context = compact_workflow_atom_context(
        build_workflow_atom_context(
            user_prompt=user_prompt or title,
            style_preset=style_preset,
            slide_count=8,
            include_prompt=False,
        )
    )
    deck_style: dict[str, Any] = {
        "visual_density": "medium",
        "emoji_mode": "none",
    }
    for key, value in (atom_context.get("deck_style_delta") or {}).items():
        if key in {"chart_treatment", "table_treatment", "header_mode", "footer_mode", "visual_density"}:
            deck_style[key] = value
    if font_pair:
        deck_style["font_pair"] = font_pair
    if palette_key:
        deck_style["palette_key"] = palette_key
    return {
        "title": title,
        "subtitle": "Working outline",
        "deck_style": deck_style,
        "metadata": {
            "starter_outline_version": "style_reference_starter_outline_v1",
            "starter_outline_status": "synthetic_scaffold_replace_before_delivery",
            "style_reference": {
                "catalog_version": reference.get("catalog_version"),
                "reference_id": reference.get("reference_id"),
                "reference_name": reference.get("reference_name"),
                "playbook_version": playbook.get("playbook_version"),
                "style_metric_profile": reference.get("style_metric_profile"),
            },
            "style_atom_context": atom_context,
        },
        "slides": [
            {
                "slide_id": "s1",
                "type": "title",
                "title": title,
                "subtitle": str(reference.get("style_dna") or "Prepare notes, assets, and outline before building"),
                "chips": [
                    _compact_text(style_preset, "style preset", limit=18),
                    _compact_text(reference.get("reference_name"), "style ref", limit=18),
                    "playbook-v1" if playbook.get("playbook_version") == LAYOUT_PLAYBOOK_VERSION else str(playbook.get("playbook_version") or "playbook"),
                ],
            },
            {
                "slide_id": "s2",
                "type": "content",
                "variant": "split",
                "title": "Core message",
                "subtitle": "Start from the decision or takeaway",
                "bullets": [
                    "State the main decision, result, or recommendation first.",
                    "Use one evidence object per content slide when data is available.",
                    "Keep source-backed claims explicit and traceable to the planning files.",
                    "Convert dense prose into a chart, table, figure, or short comparison.",
                ],
                "highlights": [
                    "Alignment and readability are release gates.",
                    "Assets stay source-backed and optional.",
                    "QA blocks overlap, overflow, and sparse layouts.",
                ],
            },
            *_style_reference_starter_slides(style_preset),
        ],
    }


def _write_text(path: Path, content: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content, encoding="utf-8")


def _hex_rgb(value: Any, fallback: tuple[int, int, int]) -> tuple[int, int, int]:
    text = str(value or "").strip().lstrip("#")
    if len(text) != 6:
        return fallback
    try:
        return (int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16))
    except ValueError:
        return fallback


def _blend_rgb(a: tuple[int, int, int], b: tuple[int, int, int], t: float) -> tuple[int, int, int]:
    return tuple(max(0, min(255, int(round(x + (y - x) * t)))) for x, y in zip(a, b))


def _png_chunk(kind: bytes, payload: bytes) -> bytes:
    return (
        struct.pack(">I", len(payload))
        + kind
        + payload
        + struct.pack(">I", zlib.crc32(kind + payload) & 0xFFFFFFFF)
    )


def _write_synthetic_png(
    path: Path,
    *,
    bg: tuple[int, int, int],
    surface: tuple[int, int, int],
    accent: tuple[int, int, int],
    accent_2: tuple[int, int, int],
    width: int = 960,
    height: int = 540,
) -> None:
    pixels = bytearray(bg * width * height)

    def rect(x0: int, y0: int, x1: int, y1: int, color: tuple[int, int, int]) -> None:
        x0c, y0c = max(0, x0), max(0, y0)
        x1c, y1c = min(width, x1), min(height, y1)
        for y in range(y0c, y1c):
            row = y * width * 3
            for x in range(x0c, x1c):
                offset = row + x * 3
                pixels[offset : offset + 3] = bytes(color)

    rect(38, 36, width - 38, height - 36, surface)
    rect(70, 74, width - 70, 82, accent)
    rect(70, height - 92, width - 70, height - 86, _blend_rgb(accent, surface, 0.42))
    panel_w = (width - 180) // 3
    for idx in range(3):
        x = 70 + idx * (panel_w + 20)
        y = 120
        rect(x, y, x + panel_w, y + 250, _blend_rgb(surface, bg, 0.16))
        rect(x, y, x + panel_w, y + 5, accent if idx % 2 == 0 else accent_2)
        for bar_idx, value in enumerate((0.42, 0.68, 0.54, 0.78)):
            bar_x = x + 28 + bar_idx * 44
            bar_h = int(160 * value)
            rect(bar_x, y + 212 - bar_h, bar_x + 25, y + 212, accent if bar_idx % 2 == 0 else accent_2)
        rect(x + 28, y + 224, x + panel_w - 28, y + 229, _blend_rgb(bg, accent, 0.35))
    for idx, value in enumerate((0.74, 0.56, 0.83, 0.65)):
        x = 86 + idx * 205
        rect(x, 418, x + int(120 * value), 434, accent if idx % 2 == 0 else accent_2)
        rect(x, 446, x + 148, 452, _blend_rgb(bg, surface, 0.42))

    raw = b"".join(
        b"\x00" + bytes(pixels[y * width * 3 : (y + 1) * width * 3])
        for y in range(height)
    )
    payload = (
        b"\x89PNG\r\n\x1a\n"
        + _png_chunk(b"IHDR", struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0))
        + _png_chunk(b"IDAT", zlib.compress(raw, 9))
        + _png_chunk(b"IEND", b"")
    )
    path.parent.mkdir(parents=True, exist_ok=True)
    if not path.exists() or path.read_bytes() != payload:
        path.write_bytes(payload)


def _write_style_reference_assets(workspace: Path, style_preset: str) -> dict[str, str]:
    preset = PRESETS[style_preset]
    palette = preset.palette
    bg = _hex_rgb(palette.get("bg"), (248, 250, 252))
    surface = _hex_rgb(palette.get("surface"), (255, 255, 255))
    accent = _hex_rgb(palette.get("accent_primary"), (37, 99, 235))
    accent_2 = _hex_rgb(palette.get("accent_secondary"), (14, 165, 233))
    style_dir = workspace / "assets" / "style_reference"
    figure_paths = [
        style_dir / "starter_figure_a.png",
        style_dir / "starter_figure_b.png",
        style_dir / "starter_figure_c.png",
    ]
    for idx, path in enumerate(figure_paths):
        _write_synthetic_png(
            path,
            bg=bg,
            surface=surface,
            accent=accent if idx != 1 else accent_2,
            accent_2=accent_2 if idx != 1 else accent,
        )
    flow_path = style_dir / "starter_flow.mmd"
    _write_text(
        flow_path,
        "flowchart LR\n"
        "  A[Input] --> B[Check]\n"
        "  B --> C[Route]\n"
        "  C --> D[Output]\n",
    )
    _write_synthetic_png(
        workspace / STYLE_REFERENCE_ASSETS["icon"],
        bg=surface,
        surface=_blend_rgb(surface, accent, 0.12),
        accent=accent,
        accent_2=accent_2,
        width=256,
        height=256,
    )
    return {
        "figure_a": "assets/style_reference/starter_figure_a.png",
        "figure_b": "assets/style_reference/starter_figure_b.png",
        "figure_c": "assets/style_reference/starter_figure_c.png",
        "flow": "assets/style_reference/starter_flow.mmd",
        "icon": STYLE_REFERENCE_ASSETS["icon"],
    }


def _workspace_readme(slug: str, title: str) -> str:
    return f"""# {title}

This workspace is the saved authoring source for the `{slug}` deck.

## Files

- `outline.json`: canonical structured slide source
- `content_plan.json`: thesis, audience, slide roles, and visual strategy
- `design_brief.json`: audience posture, cover concept, structure strategy, and grid policy
- `evidence_plan.json`: sourced claims, metrics, chart candidates, and gaps
- `style_contract.json`: stable style + layout contract for later slide additions
- `asset_plan.json`: source-backed imagery/background/chart staging plan
- `notes.md`: deck-specific data sources, decisions, and manual design notes
- `data/`: local datasets copied or linked for reproducible analysis
- `assets/data/`: smaller data extracts or tables staged with the deck
- `assets/figures/`: generated slide-ready figures
- `assets/charts/`: generated editable chart JSON specs
- `assets/`: local images, diagrams, logos, and tables used by the deck
- `build/`: generated `.pptx` output plus QA reports

## Commands

Emit the first-turn packet during initialization for reproducible intake and
design-contract handoff:

```bash
python3 ../../scripts/init_deck_workspace.py --workspace . --title "{title}" --style-preset <preset> --user-prompt "Original user request"
```

Build the deck:

```bash
python3 ../../scripts/build_workspace.py --workspace . --overwrite
```

Build and run strict QA:

```bash
python3 ../../scripts/build_workspace.py --workspace . --qa --overwrite
```

Final reusable/report build with planning and whitespace gates:

```bash
python3 ../../scripts/build_workspace.py --workspace . --qa --fail-on-planning-warnings --fail-on-whitespace-warnings --overwrite
```

Use non-render QA when LibreOffice is unavailable:

```bash
python3 ../../scripts/build_workspace.py --workspace . --qa --skip-render --overwrite
```

Build, scaffold local data artifacts, and run QA:

```bash
python3 ../../scripts/build_workspace.py --workspace . --scaffold-data-artifacts --auto-bind-artifacts --qa --fail-on-planning-warnings --fail-on-whitespace-warnings --overwrite
```

Fail final polish on accidental dead whitespace:

```bash
python3 ../../scripts/build_workspace.py --workspace . --qa --fail-on-whitespace-warnings --overwrite
```

Allow Wikimedia Commons fetches while staging assets:

```bash
python3 ../../scripts/build_workspace.py --workspace . --allow-network-assets --overwrite
```

## Iteration Pattern

1. Fill `content_plan.json` with thesis, audience, slide roles, and visual strategy.
2. Fill `design_brief.json` with audience posture, cover concept, and structure strategy.
3. Fill `evidence_plan.json` with sourced claims, metrics, and chart candidates.
4. Update `notes.md` with data rules and unresolved assumptions.
5. Add source-backed image/background/chart requests to `asset_plan.json`.
6. Put local CSV/TSV/XLSX/JSON data in `data/` or `assets/data/`.
7. Run `python3 ../../scripts/build_workspace.py --workspace . --scaffold-data-artifacts --auto-bind-artifacts --qa --fail-on-planning-warnings --fail-on-whitespace-warnings --overwrite`
   when a dataset should become a repeatable chart/figure artifact as part of the build, or run
   `python3 ../../scripts/scaffold_figure_artifacts.py --workspace . --run`
   when you want a separate scaffold/refine step.
8. Stage local assets inside `assets/` when needed.
9. Edit `outline.json` to add, replace, or reorder slides.
10. Reference staged assets with aliases such as `asset:hero_name`, `image:crew_portrait`, `chart:result_chart`, or `generated:concept_visual`.
11. Re-run `build_workspace.py`.
12. Before final delivery, run `python3 ../../scripts/build_workspace.py --workspace . --qa --fail-on-planning-warnings --fail-on-whitespace-warnings --overwrite`.
13. Keep the source files. Do not rely on inline heredoc generation if you want to extend the deck later.
"""


def _workspace_notes(title: str, style_preset: str) -> str:
    preset = PRESETS[style_preset]
    style_reference = preset_style_reference(style_preset)
    playbook = _story_dict(style_reference.get("layout_playbook"))
    metric_profile = _story_dict(style_reference.get("style_metric_profile"))
    body_budget = metric_profile.get("body_words_per_content_slide")
    body_budget_text = ", ".join(str(item) for item in body_budget) if isinstance(body_budget, list) else ""
    content_treatments = _story_dict(style_reference.get("content_treatments"))
    return f"""# {title} Notes

## Purpose

- Audience:
- Decision / outcome:
- Style preset: `{style_preset}`
- Style reference: `{style_reference.get("reference_id")}` / {style_reference.get("reference_name")}
- Style metrics: `{metric_profile.get("metric_profile_version")}`; density `{metric_profile.get("density_level")}`; whitespace target `{metric_profile.get("whitespace_ratio_target")}`; body-word budget `{body_budget_text}`.
- Starter scaffold: `style_reference_starter_outline_v1` synthetic examples; replace before delivery.

## Sources

- Add the datasets, URLs, or reference decks used to author this presentation.
- Record the provenance for every non-user image you stage through `asset_plan.json`.
- Promote researched claims into `evidence_plan.json` before adding them to slides.

## Research log to staging plan

Closes the gap where research produces good content but never turns into
staged visuals. Every row in this table should eventually trigger an
entry in `asset_plan.json` (wikimedia_query for a CC photo, or a staged
icon/chart).

| Fact discovered | Source | Becomes | In asset_plan as |
|---|---|---|---|
| _e.g. Chicago Pile-1, first controlled chain reaction, Dec 2 1942_ | _en.wikipedia.org/Chicago_Pile-1_ | _hero image on slide 3_ | _images[0].wikimedia_query: "Chicago Pile-1"_ |
|  |  |  |  |
|  |  |  |  |

If this table is empty at build time, ask yourself whether the deck
actually has no visual anchors or whether the research hasn't been
connected to the staging plan yet.

## Style Contract

- Slide size: 16:9 unless a reference deck says otherwise
- Title font: {preset.typography.title_max}-{preset.typography.title_min}pt range via preset
- Section font: {preset.typography.section_max}-{preset.typography.section_min}pt range via preset
- Body font: {preset.typography.body_max}-{preset.typography.body_min}pt range via preset
- Margin x: {preset.layout.margin_x}
- Gutter: {preset.layout.gutter}
- Style DNA: {style_reference.get("style_dna")}
- Preferred variants: {", ".join(str(item) for item in _story_list(playbook.get("preferred_variants"))[:8])}
- Chart treatment: {content_treatments.get("chart", "")}
- Table treatment: {content_treatments.get("table", "")}
- Decision treatment: {content_treatments.get("decision", "")}

## QA Notes

- Preserve alignment first.
- Keep subtitles below wrapped titles.
- Prefer local, source-backed assets in `assets/`.
- Use `asset:alias` references in `outline.json` after staging into `assets/staged/`.
- Add any deck-specific measurements here if you later match an existing deck manually.
"""


def _resolve_workspace_output(workspace: Path, raw: str) -> Path:
    path = Path(str(raw or "")).expanduser()
    if path.is_absolute():
        return path.resolve()
    return (workspace / path).resolve()


def _display_path(workspace: Path, path: Path) -> str:
    try:
        return str(path.resolve().relative_to(workspace.resolve()))
    except ValueError:
        return str(path.resolve())


def _slide_role_from_ref(ref: dict[str, str], index: int) -> str:
    variant = str(ref.get("variant") or "").strip().lower()
    treatment = str(ref.get("treatment_key") or "").strip().lower()
    if variant == "title" or index == 0:
        return "title"
    if treatment in {"chart", "table", "figure", "dashboard"}:
        return "evidence"
    if treatment == "comparison" or variant in {"comparison-2col", "split", "matrix"}:
        return "comparison"
    if variant in {"timeline", "flow"}:
        return "sequence"
    if treatment == "decision" or variant == "standard":
        return "decision"
    return "setup"


def _visual_strategy_from_ref(ref: dict[str, str]) -> str:
    variant = str(ref.get("variant") or "").strip().lower()
    strategies = {
        "title": "title opener with preset-specific style reference chips",
        "split": "structured two-column contrast or setup",
        "chart": "editable chart scaffold from the preset's synthetic storyboard",
        "table": "editable table scaffold from the preset's content-treatment grammar",
        "lab-run-results": "compact lab/report results table scaffold",
        "comparison-2col": "two-column comparison with visible verdict",
        "stats": "dashboard fact tiles with short readouts",
        "kpi-hero": "single hero metric rhythm break",
        "timeline": "milestone sequence only when order matters",
        "matrix": "2x2 decision or tradeoff matrix",
        "cards-2": "two parallel proof blocks only when content is genuinely paired",
        "cards-3": "parallel proof blocks, not default decoration",
        "image-sidebar": "local starter figure plus interpretation sidebar",
        "scientific-figure": "multi-panel starter figure grid with caption and interpretation",
        "flow": "local Mermaid process diagram with sidebar interpretation",
    }
    return strategies.get(variant, "preset-specific evidence or synthesis scaffold")


def _content_plan_stub(title: str, slide_refs: list[dict[str, str]] | None = None) -> dict[str, Any]:
    refs = slide_refs or [{"slide_id": "s1", "variant": "title"}, {"slide_id": "s2", "variant": "split"}]
    first_slide_id = refs[0].get("slide_id") or "s1"
    setup_ids = [ref.get("slide_id") or f"s{index + 1}" for index, ref in enumerate(refs[:2])]
    evidence_ids = [
        ref.get("slide_id") or f"s{index + 1}"
        for index, ref in enumerate(refs)
        if index >= 2 and _slide_role_from_ref(ref, index) in {"evidence", "comparison", "sequence"}
    ]
    implication_ids = [
        ref.get("slide_id") or f"s{index + 1}"
        for index, ref in enumerate(refs)
        if index >= 2 and _slide_role_from_ref(ref, index) == "decision"
    ]
    if not implication_ids and refs:
        last_id = refs[-1].get("slide_id") or f"s{len(refs)}"
        if last_id not in evidence_ids and last_id not in setup_ids:
            implication_ids = [last_id]
    slide_plan = []
    for index, ref in enumerate(refs):
        slide_id = ref.get("slide_id") or f"s{index + 1}"
        variant = ref.get("variant") or ("title" if index == 0 else "standard")
        role = _slide_role_from_ref(ref, index)
        title_text = ref.get("title") or ("Open the deck" if index == 0 else "Replace scaffold content")
        starter_kind = ref.get("starter_kind") or ""
        slide_plan.append(
            {
                "slide_id": slide_id,
                "role": role,
                "message": (
                    "Replace this synthetic style-reference scaffold with topic-specific content."
                    if starter_kind == "style_reference"
                    else (
                        "Start from a durable workspace, not one-off inline code."
                        if index == 0
                        else "The durable source files are the contract for future edits."
                    )
                ),
                "variant": variant,
                "visual_strategy": _visual_strategy_from_ref(ref),
                "evidence_needs": [] if starter_kind != "style_reference" else ["replace_synthetic_style_reference_content"],
                "asset_needs": [],
                "source_status": "synthetic_style_reference_scaffold" if starter_kind == "style_reference" else "workspace_starter",
                "outline_title": title_text,
            }
        )
    return {
        "topic": title,
        "audience": "Deck author using the presentation-skill workspace scaffold.",
        "objective": "Replace the starter content with topic-specific narrative, evidence, and assets.",
        "thesis": "A reliable deck starts with a preset-specific style reference, sourced evidence, staged visuals, and QA before delivery.",
        "narrative_arc": [
            {
                "act": "setup",
                "purpose": "Frame why the topic matters.",
                "slides": setup_ids,
            },
            {
                "act": "style-reference scaffold",
                "purpose": "Show the selected preset's reusable chart, table, comparison, dashboard, or sequence grammar before topic-specific authoring.",
                "slides": evidence_ids,
            },
            {
                "act": "implication",
                "purpose": "Close with what the audience should remember or do.",
                "slides": implication_ids,
            },
        ],
        "slide_plan": slide_plan,
        "design_notes": {
            "style_preset_reason": "Starter uses the requested preset while keeping typography and spacing conservative.",
            "rhythm_break": "Add a diagram, figure, table, or image only when the deck topic makes it useful.",
            "visual_motif": "Source-first authoring with clear stage labels when the author chooses a process visual.",
            "starter_slide_id": first_slide_id,
            "style_reference_scaffold": "Slides marked source_status=synthetic_style_reference_scaffold are style memory only and should be replaced before final delivery.",
        },
    }


def _design_brief_stub(title: str, style_preset: str, *, user_prompt: str = "") -> dict[str, Any]:
    preset = PRESETS[style_preset]
    treatment_profile = preset_treatment_profile(style_preset)
    style_reference = preset_style_reference(style_preset)
    playbook = _story_dict(style_reference.get("layout_playbook"))
    atom_context = compact_workflow_atom_context(
        build_workflow_atom_context(
            user_prompt=user_prompt or title,
            style_preset=style_preset,
            slide_count=8,
            include_prompt=False,
        )
    )
    atom_brief = (
        atom_context.get("design_brief_delta")
        if isinstance(atom_context.get("design_brief_delta"), dict)
        else {}
    )
    style_atom_composition = (
        atom_brief.get("style_atom_composition")
        if isinstance(atom_brief.get("style_atom_composition"), dict)
        else atom_context.get("style_atom_composition")
    )
    preferred_variants = [
        str(item)
        for item in _story_list(playbook.get("preferred_variants"))
        if str(item).strip()
    ]
    return {
        "topic": title,
        "content_maturity": "serious/work",
        "audience_posture": "coworkers/operators",
        "emotional_register": "trustworthy",
        "format_promise": (
            "A clean, editable PowerPoint deck with one dominant idea per slide, "
            "disciplined alignment, and enough visual rhythm to avoid generic card grids."
        ),
        "anti_format": [
            "repeated title plus three cards on every slide",
            "body text placed by feel instead of grid constants",
            "decorative shapes without a reading job",
            "shrinking text below readability floors to solve density",
        ],
        "canvas_and_grid": {
            "aspect": "16:9",
            "margin_x_in": 0.5,
            "footer_reserve_in": 0.32,
            "header_policy": "measured title/subtitle stack; body starts at returned contentTop",
            "column_policy": "derive columns from margin and gutter constants, not magic numbers",
        },
        "visual_system": {
            "style_preset": style_preset,
            "dominant_color": preset.palette["bg_dark"],
            "accent_primary": preset.palette["accent_primary"],
            "accent_secondary": preset.palette["accent_secondary"],
            "palette_role_map": {
                "background": "dominant or neutral field",
                "accent": "navigation, KPI emphasis, rails, and labels",
                "muted": "captions and provenance",
            },
        },
        "style_system": {
            "style_preset": style_preset,
            "style_seed": f"{_slugify(title)}-{style_preset}",
            "preset_treatment_profile": treatment_profile,
            "style_reference": style_reference,
            "style_atom_context": atom_context,
            "style_atom_composition": style_atom_composition,
            "style_atom_preferred_variants": atom_context.get("preferred_variants") or [],
            "style_atom_narrative_arc": atom_context.get("narrative_arc") or [],
            "style_mix_matrix": treatment_profile["style_mix_matrix"],
        },
        "style_atom_composition": style_atom_composition,
        "palette_signals": atom_brief.get("palette_signals", []),
        "typography_signals": atom_brief.get("typography_signals", []),
        "layout_signals": atom_brief.get("layout_signals", []),
        "rhythm_signature": atom_brief.get("rhythm_signature", ""),
        "title_page_concept": {
            "chosen_archetype": "topic-specific opener chosen from the preset and content",
            "dominant_element": "large topic-specific title",
            "supporting_element": "short subtitle or one strong hero asset",
            "why_this_could_only_be_this_deck": "Replace with a sentence before final delivery.",
        },
        "structure_strategy": {
            "primary_scaffold": "open editorial content slides with measured headers",
            "repeated_elements": ["shared margins", "consistent source/footer treatment", "limited accent rails"],
            "allowed_variations": preferred_variants
            or [
                "standard clean report slides",
                "split",
                "cards-2",
                "cards-3 with promote_card",
                "timeline only when the sequence is truly time-based",
                "table",
                "matrix",
                "optional kpi-hero only when one metric deserves isolation",
                "flow",
                "generated-image",
            ],
            "style_reference_layout_playbook": {
                "playbook_version": playbook.get("playbook_version"),
                "reference_id": playbook.get("reference_id"),
                "opening_sequence": playbook.get("opening_sequence"),
                "content_rules": playbook.get("content_rules"),
                "avoid_variants": playbook.get("avoid_variants"),
            },
            "container_policy": (
                "Cards are for modular comparisons or evidence groups, not the default "
                "way to make prose look designed."
            ),
            "rhythm_break_plan": (
                "Use a rhythm break only when the content asks for it: a true "
                "hero metric, a full-bleed/source-backed image, a major section "
                "turn, or a decisive chart. Do not add a KPI hero just to break "
                "rhythm. Do not add a timeline just because a slide has steps; "
                "use report bands, a table, or a figure when those are clearer."
            ),
        },
        "readability_contract": {
            "min_title_pt": 24,
            "min_body_pt": 12,
            "min_caption_pt": 7.5,
            "max_title_lines": 2,
            "max_slide_text_lines": 8,
            "max_slide_words": 105,
            "max_slide_chars": 700,
            "footer_reserved_inches": 0.34,
            "chart_label_min_pt": 8,
            "table_density_rule": "split or summarize tables that force unreadable text",
            "whitespace_rule": "avoid awkward empty regions; choose variants that fit actual evidence shape",
            "figure_crop_rule": "tight bounding boxes and trimmed exterior whitespace",
        },
        "speed_contract": {
            "renderer": "pptxgenjs by default; Python fallback only for legacy renderer-specific behavior",
            "first_pass": "run validate_planning.py, preflight.py, and render-free QA before slide rendering",
            "render_policy": "render only after source files are stable or when visual judgment matters",
            "asset_policy": "reuse local/generated artifacts before network assets unless source-backed imagery is needed",
            "conversion_hint": "use persistent LibreOffice/unoserver when available for repeated render QA",
        },
    }


def _evidence_plan_stub(title: str) -> dict[str, Any]:
    return {
        "topic": title,
        "source_policy": "Prefer primary or source-backed facts. Do not fabricate citations.",
        "items": [],
        "chart_candidates": [],
        "open_questions": [
            "Replace the scaffold with topic-specific evidence before delivering a factual deck."
        ],
    }


def _asset_plan_stub(title: str) -> dict[str, Any]:
    """Starter plan for staged deck assets.

    The entries below are TODO placeholders (empty arrays with inline
    schema comments). They're NOT ready-to-run examples, because
    generic examples tend to ship unchanged.

    Populate the arrays with topic-specific requests, or delete the
    file if the deck doesn't need staged assets. `build_workspace.py`
    warns at build time if this file is still at its initial state and
    the deck has no icons, hero image, or charts anywhere in its
    outline.
    """
    return {
        "topic": title,
        "__readme__": (
            "Delete this __readme__ key and populate the arrays below "
            "with real image/chart/icon requests for THIS topic. "
            "See references/outline_schema.md and "
            "references/deck_workspace_mode.md for the schemas. If the "
            "deck doesn't need staged assets, delete this whole file."
        ),
        "images": [
            # Example schema - delete and replace with real entries:
            # {"name": "hero_photo", "wikimedia_query": "<topic keyword>",
            #  "allow_sharealike": true, "attribution_line": "<caption>"}
        ],
        "backgrounds": [
            # Example schema - delete and replace:
            # {"name": "section_bg", "path": "assets/staged/section_bg.png"}
        ],
        "charts": [
            # Example schema - delete and replace with real data:
            # {"name": "trend_by_year", "title": "<chart title>",
            #  "type": "line" | "bar" | "pie",
            #  "series": [{"name": "<series>", "labels": [...], "values": [...]}],
            #  "options": {"catAxisTitle": "...", "valAxisTitle": "..."}}
        ],
        "tables": [
            # Example schema - delete and replace with real data:
            # {"name": "run_summary", "path": "assets/tables/run_summary.json"}
            # where the JSON contains {"headers": [...], "rows": [[...]], "caption": "..."}.
        ],
        "generated_images": [
            # Example schema - delete and replace with deliberate generated
            # concept art only when source-backed imagery is insufficient:
            # {"name": "concept_visual",
            #  "prompt": "A precise editorial illustration of ...",
            #  "purpose": "Optional visual anchor slide",
            #  "model": "gpt-image-2",
            #  "size": "1536x1024",
            #  "quality": "medium"}
        ],
        "icons": [
            # Example schema - delete and replace with real icon names.
            # Bare names resolve against <workspace>/assets/icons/<name>.png.
            # Use when a cards-3 / timeline / matrix / stats / cards-2
            # slide has a clear visual metaphor per card.
            # {"name": "reactor_core", "path": "assets/icons/reactor_core.png"}
        ],
    }


def _style_contract(
    *,
    title: str,
    slug: str,
    style_preset: str,
    font_pair: str | None,
    palette_key: str | None,
    reference_pptx: Path | None,
    user_prompt: str = "",
) -> dict[str, Any]:
    preset = PRESETS[style_preset]
    style_reference = preset_style_reference(style_preset)
    playbook = _story_dict(style_reference.get("layout_playbook"))
    atom_context = compact_workflow_atom_context(
        build_workflow_atom_context(
            user_prompt=user_prompt or title,
            style_preset=style_preset,
            slide_count=8,
            include_prompt=False,
        )
    )
    contract: dict[str, Any] = {
        "workspace_version": 1,
        "deck_title": title,
        "deck_slug": slug,
        "build": {
            "style_preset": style_preset,
            "font_pair": font_pair,
            "palette_key": palette_key,
            "output_pptx": f"build/{slug}.pptx",
            "qa_dir": "build/qa",
            "qa_report": "build/qa/report.json",
        },
        "layout_rules": {
            "alignment_first": True,
            "zero_overlap_required": True,
            "title_subtitle_stack_dynamic": True,
            "footer_safe_region_required": True,
            "cards_use_measured_text_fit": True,
        },
        "preset_tokens": preset.to_dict(),
        "style_reference": {
            "catalog_version": style_reference.get("catalog_version"),
            "reference_id": style_reference.get("reference_id"),
            "reference_name": style_reference.get("reference_name"),
            "source_status": style_reference.get("source_status"),
            "style_dna": style_reference.get("style_dna"),
            "style_metric_profile": style_reference.get("style_metric_profile"),
            "layout_playbook_version": playbook.get("playbook_version"),
            "preferred_variants": playbook.get("preferred_variants"),
            "starter_outline_version": "style_reference_starter_outline_v1",
        },
        "style_atom_context": atom_context,
    }
    if reference_pptx:
        contract["reference"] = _reference_summary(reference_pptx)
    return contract


def _args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Initialize a persistent PPTX deck workspace.")
    parser.add_argument("--workspace", required=True, help="Workspace directory to create")
    parser.add_argument("--title", required=True, help="Human-readable deck title")
    parser.add_argument("--style-preset", default="executive-clinical", choices=sorted(PRESETS))
    parser.add_argument("--font-pair", help="Optional font pair override stored in outline/style contract")
    parser.add_argument("--palette-key", help="Optional palette override stored in outline/style contract")
    parser.add_argument("--source-outline", help="Optional JSON outline to copy into workspace")
    parser.add_argument("--reference-pptx", help="Optional reference deck to summarize and extract")
    parser.add_argument("--overwrite", action="store_true", help="Replace existing workspace files")
    parser.add_argument(
        "--user-prompt",
        default="",
        help=(
            "Original user request. When provided, init also writes a "
            "reproducible deck-start packet for intake/design-contract handoff."
        ),
    )
    parser.add_argument(
        "--emit-start-packet",
        action="store_true",
        help=(
            "Write deck_start_packet.json during initialization. If "
            "--user-prompt is omitted, the deck title is used as the prompt."
        ),
    )
    parser.add_argument(
        "--start-packet",
        default="deck_start_packet.json",
        help="Workspace-relative or absolute path for the optional deck-start packet.",
    )
    parser.add_argument(
        "--start-packet-mode",
        choices=["concise", "full"],
        default="concise",
        help="Question set size for the optional deck-start packet.",
    )
    parser.add_argument(
        "--agent-profile",
        choices=sorted(PROFILE_ALIASES),
        default="auto",
        help=(
            "Model-side execution profile for the compact agent brief. "
            "Use auto, quality-first/sol, balanced/terra, or fast/luna."
        ),
    )
    parser.add_argument(
        "--followup-edit",
        action="store_true",
        help=(
            "Explicitly acknowledge that --source-outline / --reference-pptx "
            "points at an existing deck and this is a followup edit to the "
            "SAME topic (not a new deck). Required when sourcing from another "
            "decks/<slug>/ directory. Prevents the 'clone-an-existing-deck-"
            "as-house-style' anti-pattern - see references/codex_guardrails.md."
        ),
    )
    return parser.parse_args()


def _is_under_decks_dir(path: Path) -> tuple[bool, str | None]:
    """Return (True, sibling_slug) if `path` lives under a decks/<slug>/ tree
    that is NOT the workspace being initialized. sibling_slug is the name
    of the source workspace, used for the error message.
    """
    resolved = path.resolve()
    for ancestor in resolved.parents:
        if ancestor.name == "decks" and ancestor.parent.name in {"presentation-skill", "pptx-skill"}:
            # Immediate child of decks/ is the source workspace slug.
            try:
                rel = resolved.relative_to(ancestor)
                return True, rel.parts[0] if rel.parts else None
            except ValueError:
                return False, None
    return False, None


def main() -> int:
    args = _args()
    workspace = Path(args.workspace).expanduser().resolve()
    source_outline = Path(args.source_outline).expanduser().resolve() if args.source_outline else None
    reference_pptx = Path(args.reference_pptx).expanduser().resolve() if args.reference_pptx else None
    user_prompt = str(args.user_prompt or "").strip()
    emit_start_packet = bool(args.emit_start_packet or user_prompt)
    start_packet_path = _resolve_workspace_output(workspace, args.start_packet)

    if source_outline and not source_outline.exists():
        raise FileNotFoundError(f"Source outline not found: {source_outline}")
    if reference_pptx and not reference_pptx.exists():
        raise FileNotFoundError(f"Reference deck not found: {reference_pptx}")
    if workspace.exists() and any(workspace.iterdir()) and not args.overwrite:
        raise FileExistsError(f"Workspace already exists and is not empty: {workspace}")

    # Guardrail: cloning an existing deck workspace as a "house style" for a
    # new topic is a documented Codex anti-pattern (see codex_guardrails.md
    # "Eighth Trap"). If --source-outline or --reference-pptx points into
    # another decks/<slug>/ tree, require --followup-edit to acknowledge
    # that the intent is editing the same topic, not cloning style.
    for source_label, source in (
        ("--source-outline", source_outline),
        ("--reference-pptx", reference_pptx),
    ):
        if source is None:
            continue
        under_decks, source_slug = _is_under_decks_dir(source)
        if not under_decks:
            continue
        workspace_slug = _slugify(args.title)
        # If the source is in the SAME workspace we're re-initializing
        # (e.g., overwriting), let it pass. The anti-pattern is
        # cross-topic cloning.
        if source_slug == workspace_slug or source_slug == workspace.name:
            continue
        if not args.followup_edit:
            # Softened from hard ERROR to a warning: studying a past deck
            # for file shape or vocabulary is legitimate; cloning its
            # variant mix wholesale is not. Trust the author.
            print(
                f"[init_deck_workspace] WARNING: {source_label} points at "
                f"workspace {source_slug!r}. If you're just reading it for "
                "file shape, fine. If you're cloning its variant mix for a "
                "new topic, reconsider - see codex_guardrails.md on uniform-"
                "deck syndrome. Pass --followup-edit to silence this warning.",
                file=sys.stderr,
            )

    workspace.mkdir(parents=True, exist_ok=True)
    for subdir in (
        "assets",
        "assets/charts",
        "assets/data",
        "assets/diagrams",
        "assets/figures",
        "build",
        "data",
    ):
        (workspace / subdir).mkdir(parents=True, exist_ok=True)
    _write_text(workspace / "assets" / ".gitkeep", "")
    _write_text(workspace / "assets" / "charts" / ".gitkeep", "")
    _write_text(workspace / "assets" / "data" / ".gitkeep", "")
    _write_text(workspace / "assets" / "figures" / ".gitkeep", "")
    _write_text(workspace / "data" / ".gitkeep", "")
    _write_text(workspace / "build" / ".gitkeep", "")
    if not source_outline and not reference_pptx:
        _write_style_reference_assets(workspace, args.style_preset)

    slug = _slugify(args.title)
    if source_outline:
        outline = _copy_json(source_outline)
    elif reference_pptx:
        outline = _extract_outline(reference_pptx)
    else:
        outline = _starter_outline(
            args.title,
            args.style_preset,
            args.font_pair,
            args.palette_key,
            user_prompt=user_prompt,
        )

    slide_refs = _ensure_outline_slide_ids(outline)
    outline.setdefault("title", args.title)
    if args.font_pair or args.palette_key:
        deck_style = outline.setdefault("deck_style", {})
        if args.font_pair:
            deck_style.setdefault("font_pair", args.font_pair)
        if args.palette_key:
            deck_style.setdefault("palette_key", args.palette_key)

    style_contract = _style_contract(
        title=args.title,
        slug=slug,
        style_preset=args.style_preset,
        font_pair=args.font_pair,
        palette_key=args.palette_key,
        reference_pptx=reference_pptx,
        user_prompt=user_prompt,
    )
    workspace_manifest = {
        "workspace_version": 1,
        "deck_title": args.title,
        "deck_slug": slug,
        "style_contract": "style_contract.json",
        "content_plan": "content_plan.json",
        "design_brief": "design_brief.json",
        "evidence_plan": "evidence_plan.json",
        "outline": "outline.json",
        "asset_plan": "asset_plan.json",
        "notes": "notes.md",
        "assets_dir": "assets",
        "staged_assets_dir": "assets/staged",
        "build_dir": "build",
        "reference_pptx": str(reference_pptx) if reference_pptx else None,
    }
    if emit_start_packet:
        workspace_manifest["deck_start_packet"] = _display_path(workspace, start_packet_path)
        workspace_manifest["agent_brief"] = "agent_brief.json"
        workspace_manifest["agent_brief_markdown"] = "agent_brief.md"

    _write_text(workspace / "outline.json", json.dumps(outline, indent=2, ensure_ascii=False) + "\n")
    _write_text(
        workspace / "style_contract.json",
        json.dumps(style_contract, indent=2, ensure_ascii=False) + "\n",
    )
    _write_text(
        workspace / "content_plan.json",
        json.dumps(_content_plan_stub(args.title, slide_refs), indent=2, ensure_ascii=False) + "\n",
    )
    _write_text(
        workspace / "design_brief.json",
        json.dumps(
            _design_brief_stub(args.title, args.style_preset, user_prompt=user_prompt),
            indent=2,
            ensure_ascii=False,
        )
        + "\n",
    )
    _write_text(
        workspace / "evidence_plan.json",
        json.dumps(_evidence_plan_stub(args.title), indent=2, ensure_ascii=False) + "\n",
    )
    _write_text(workspace / "asset_plan.json", json.dumps(_asset_plan_stub(args.title), indent=2) + "\n")
    _write_text(workspace / "workspace.json", json.dumps(workspace_manifest, indent=2) + "\n")
    _write_text(workspace / "README.md", _workspace_readme(slug, args.title))
    _write_text(workspace / "notes.md", _workspace_notes(args.title, args.style_preset))
    if emit_start_packet:
        packet = build_packet(
            workspace=workspace,
            user_prompt=user_prompt or args.title,
            mode=args.start_packet_mode,
        )
        _write_text(
            start_packet_path,
            json.dumps(packet, indent=2, ensure_ascii=False) + "\n",
        )
        write_agent_brief(
            packet=packet,
            workspace=workspace,
            user_prompt=user_prompt or args.title,
            requested_profile=args.agent_profile,
        )

    print(f"Workspace created: {workspace}")
    print(f"Outline: {workspace / 'outline.json'}")
    print(f"Style contract: {workspace / 'style_contract.json'}")
    if emit_start_packet:
        print(f"Deck start packet: {start_packet_path}")
        print(f"Agent brief: {workspace / 'agent_brief.md'}")
    print(f"Build target: {workspace / 'build' / f'{slug}.pptx'}")
    return 0


if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except Exception as exc:  # pragma: no cover - CLI error path
        print(f"Error: {exc}")
        raise SystemExit(1)
